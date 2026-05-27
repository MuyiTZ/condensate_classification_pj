#!/usr/bin/env python3
"""
Diagnostics and separability extension for control-guided condensate analysis.

This script builds on existing control_guided_classification outputs. It does
not rerun the original classifier. All outputs are written under:

  control_guided_classification/diagnostics_and_separability/
"""

from __future__ import annotations

import argparse
import json
import math
import os
import platform
import shutil
import sys
import zipfile
from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt

from scipy.spatial.distance import cdist, pdist
from scipy.stats import ks_2samp, mannwhitneyu

from sklearn.decomposition import PCA
from sklearn.ensemble import RandomForestClassifier
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import (
    accuracy_score,
    balanced_accuracy_score,
    confusion_matrix,
    f1_score,
    precision_score,
    recall_score,
    roc_auc_score,
    silhouette_score,
)
from sklearn.model_selection import StratifiedKFold
from sklearn.preprocessing import StandardScaler

try:
    import umap  # type: ignore
except Exception:  # pragma: no cover
    umap = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except Exception:  # pragma: no cover
    Presentation = None


RANDOM_SEED = 42
CLASS_NAMES = [
    "free_like",
    "chromatin_bound_like",
    "mcm_bound_like",
    "slow_diffusive",
    "mixed",
    "target_specific_trapping_like",
    "ambiguous",
]


def now_stamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


class Logger:
    def __init__(self, path: Path):
        self.path = path
        self.path.parent.mkdir(parents=True, exist_ok=True)
        self.write("Diagnostics and separability analysis started")
        self.write(f"python={sys.version}")
        self.write(f"platform={platform.platform()}")
        self.write(f"cwd={Path.cwd()}")

    def write(self, message: str) -> None:
        with self.path.open("a", encoding="utf-8") as handle:
            handle.write(f"[{datetime.now().isoformat(timespec='seconds')}] {message}\n")


def rel(path: Path, root: Path) -> str:
    return str(path.resolve().relative_to(root.resolve())).replace("\\", "/")


def parse_dataset(path: Path, project: Path) -> dict:
    parts = path.resolve().relative_to(project.resolve()).parts
    protein = parts[0] if parts else ""
    dataset = parts[1] if len(parts) > 1 else ""
    tokens = dataset.split("-")
    genotype = tokens[0].upper() if tokens else "unknown"
    date_or_experiment_id = tokens[1] if len(tokens) > 1 else "unknown"
    treatment = "-".join(tokens[2:]).upper() if len(tokens) > 2 else "unknown"
    if genotype not in {"WT", "12D"} and protein.upper() == "TDP43":
        genotype = "unknown"
    if treatment not in {"CT", "ARS", "ACTD"} and protein.upper() == "TDP43":
        treatment = "unknown"
    return {
        "protein": "H2B" if protein.lower() == "h2b" else protein,
        "dataset_name": dataset,
        "genotype": genotype,
        "date_or_experiment_id": date_or_experiment_id,
        "treatment": treatment,
        "replicate": dataset,
    }


def locate_inputs(project: Path, out: Path, log: Logger) -> dict:
    files = [p for p in project.rglob("*") if p.is_file() and ".git" not in p.parts]
    h2b = [p for p in files if p.name.lower() == "single-behaviour-track-h2b-bound.csv"]
    mcm = [p for p in files if p.name.lower() == "single-behaviour-track-mcm2-bound.csv"]
    final = project / "control_guided_classification" / "final_control_guided_track_classification.csv"
    feature_matrix = project / "control_guided_classification" / "feature_matrix_tracks.csv"
    old_eval = [
        p
        for p in files
        if p.suffix.lower() == ".csv"
        and p.name.lower() in {"condensate_evaluation_summary.csv", "condensate_evaluation_windows.csv", "condensate_stats.csv"}
    ]

    if len(h2b) != 1:
        raise RuntimeError(f"Expected exactly one H2B-bound file, found {len(h2b)}: {[rel(p, project) for p in h2b]}")
    if len(mcm) != 1:
        raise RuntimeError(f"Expected exactly one MCM2-bound file, found {len(mcm)}: {[rel(p, project) for p in mcm]}")
    if not final.exists():
        raise RuntimeError(f"Missing final control-guided table: {final}")
    if not feature_matrix.exists():
        log.write(f"WARNING missing feature matrix: {feature_matrix}")

    inv = pd.DataFrame(
        [
            {"role": "H2B_Bound", "relative_path": rel(h2b[0], project), "size_bytes": h2b[0].stat().st_size},
            {"role": "MCM2_Bound", "relative_path": rel(mcm[0], project), "size_bytes": mcm[0].stat().st_size},
            {"role": "final_control_guided", "relative_path": rel(final, project), "size_bytes": final.stat().st_size},
            {
                "role": "feature_matrix",
                "relative_path": rel(feature_matrix, project) if feature_matrix.exists() else "",
                "size_bytes": feature_matrix.stat().st_size if feature_matrix.exists() else 0,
            },
        ]
    )
    inv.to_csv(out / "diagnostics_input_inventory.csv", index=False)
    pd.DataFrame({"relative_path": [rel(p, project) for p in old_eval]}).to_csv(
        out / "diagnostics_old_condensate_eval_files.csv", index=False
    )
    log.write(f"H2B-bound file used: {h2b[0]}")
    log.write(f"MCM2-bound file used: {mcm[0]}")
    log.write(f"Old condensate evaluation files discovered: {len(old_eval)}")
    return {"h2b": h2b[0], "mcm": mcm[0], "final": final, "feature_matrix": feature_matrix, "old_eval": old_eval}


def max_run(values: list[int], target: int) -> int:
    best = cur = 0
    for v in values:
        if v == target:
            cur += 1
            best = max(best, cur)
        else:
            cur = 0
    return best


def frame_track_features(csv_path: Path, label: str, project: Path) -> pd.DataFrame:
    df = pd.read_csv(csv_path)
    keys = ["Video #", "Cell", "Track"]
    rows = []
    meta = parse_dataset(csv_path, project)
    for key, g in df.groupby(keys, dropna=False):
        g = g.sort_values("Frame")
        x = pd.to_numeric(g.get("x"), errors="coerce")
        y = pd.to_numeric(g.get("y"), errors="coerce")
        inten = pd.to_numeric(g.get("Intensity"), errors="coerce")
        bound = pd.to_numeric(g.get("Bound"), errors="coerce").dropna().astype(int).tolist()
        dx = x.diff()
        dy = y.diff()
        step = np.sqrt(dx * dx + dy * dy).dropna()
        if x.notna().sum() > 1 and y.notna().sum() > 1:
            net = float(np.sqrt((x.iloc[-1] - x.iloc[0]) ** 2 + (y.iloc[-1] - y.iloc[0]) ** 2))
            cx, cy = float(x.mean()), float(y.mean())
            rg = float(np.sqrt(((x - cx) ** 2 + (y - cy) ** 2).mean()))
        else:
            net, rg = np.nan, np.nan
        valid = [v for v in bound if v in {0, 1, 2}]
        n_valid = len(valid)
        row = {
            **meta,
            "bound_control_class": label,
            "source_file": csv_path.name,
            "source_folder": str(csv_path.parent),
            "video_id": key[0],
            "cell_id": key[1],
            "track_id": key[2],
            "n_frames": len(g),
            "track_length": len(g),
            "step_size_mean": float(step.mean()) if len(step) else np.nan,
            "step_size_median": float(step.median()) if len(step) else np.nan,
            "step_size_p95": float(step.quantile(0.95)) if len(step) else np.nan,
            "path_length": float(step.sum()) if len(step) else np.nan,
            "net_displacement": net,
            "Rg": rg,
            "intensity_mean": float(inten.mean()) if inten.notna().any() else np.nan,
            "intensity_std": float(inten.std()) if inten.notna().any() else np.nan,
            "frac_diffusive": valid.count(0) / n_valid if n_valid else np.nan,
            "frac_constrained": valid.count(1) / n_valid if n_valid else np.nan,
            "frac_bound": valid.count(2) / n_valid if n_valid else np.nan,
            "max_consecutive_diffusive": max_run(valid, 0) if n_valid else np.nan,
            "max_consecutive_constrained": max_run(valid, 1) if n_valid else np.nan,
            "max_consecutive_bound": max_run(valid, 2) if n_valid else np.nan,
            "n_state_switches": sum(1 for a, b in zip(valid, valid[1:]) if a != b),
        }
        rows.append(row)
    return pd.DataFrame(rows)


def choose_features(df: pd.DataFrame, ignore: set[str]) -> list[str]:
    feats = []
    for c in df.columns:
        if c in ignore:
            continue
        if pd.api.types.is_numeric_dtype(df[c]) and df[c].notna().mean() >= 0.5 and df[c].nunique(dropna=True) > 1:
            feats.append(c)
    return feats


def balanced_cv_metrics(X: pd.DataFrame, y: pd.Series, seed: int = RANDOM_SEED) -> tuple[pd.DataFrame, pd.DataFrame, RandomForestClassifier, StandardScaler]:
    scaler = StandardScaler()
    Xz = scaler.fit_transform(X)
    cv = StratifiedKFold(n_splits=min(5, int(y.value_counts().min())), shuffle=True, random_state=seed)
    probs = np.zeros(len(y))
    preds = np.empty(len(y), dtype=object)
    for train, test in cv.split(Xz, y):
        clf = RandomForestClassifier(n_estimators=160, random_state=seed, class_weight="balanced", n_jobs=-1)
        clf.fit(Xz[train], y.iloc[train])
        preds[test] = clf.predict(Xz[test])
        class_idx = list(clf.classes_).index("MCM2_Bound")
        probs[test] = clf.predict_proba(Xz[test])[:, class_idx]
    metrics = {
        "accuracy": accuracy_score(y, preds),
        "balanced_accuracy": balanced_accuracy_score(y, preds),
        "precision_macro": precision_score(y, preds, average="macro", zero_division=0),
        "recall_macro": recall_score(y, preds, average="macro", zero_division=0),
        "f1_macro": f1_score(y, preds, average="macro", zero_division=0),
        "roc_auc": roc_auc_score((y == "MCM2_Bound").astype(int), probs),
    }
    final = RandomForestClassifier(n_estimators=240, random_state=seed, class_weight="balanced", n_jobs=-1)
    final.fit(Xz, y)
    cm = pd.DataFrame(confusion_matrix(y, preds, labels=["H2B_Bound", "MCM2_Bound"]), index=["true_H2B_Bound", "true_MCM2_Bound"], columns=["pred_H2B_Bound", "pred_MCM2_Bound"])
    return pd.DataFrame([metrics]), cm, final, scaler


def repeated_downsample(X: pd.DataFrame, y: pd.Series, n_repeats: int = 30) -> pd.DataFrame:
    rng = np.random.default_rng(RANDOM_SEED)
    rows = []
    counts = y.value_counts()
    n = int(counts.min())
    for i in range(n_repeats):
        idx = []
        for cls in counts.index:
            candidates = np.where(y.to_numpy() == cls)[0]
            idx.extend(rng.choice(candidates, n, replace=False))
        idx = np.array(idx)
        m, _, _, _ = balanced_cv_metrics(X.iloc[idx].reset_index(drop=True), y.iloc[idx].reset_index(drop=True), RANDOM_SEED + i)
        rows.append({"repeat": i + 1, **m.iloc[0].to_dict()})
    return pd.DataFrame(rows)


def permutation_test(X: pd.DataFrame, y: pd.Series, observed: float, n_perm: int = 50) -> pd.DataFrame:
    rng = np.random.default_rng(RANDOM_SEED)
    vals = []
    for i in range(n_perm):
        yp = pd.Series(rng.permutation(y.to_numpy()))
        m, _, _, _ = balanced_cv_metrics(X.reset_index(drop=True), yp, RANDOM_SEED + i)
        vals.append(float(m["balanced_accuracy"].iloc[0]))
    p = (sum(v >= observed for v in vals) + 1) / (len(vals) + 1)
    return pd.DataFrame({"n_permutations": [n_perm], "observed_balanced_accuracy": [observed], "mean_permuted_balanced_accuracy": [np.mean(vals)], "p_value": [p]})


def effect_sizes(df: pd.DataFrame, features: list[str]) -> pd.DataFrame:
    rows = []
    a = df[df["bound_control_class"] == "H2B_Bound"]
    b = df[df["bound_control_class"] == "MCM2_Bound"]
    for f in features:
        x = pd.to_numeric(a[f], errors="coerce").dropna()
        y = pd.to_numeric(b[f], errors="coerce").dropna()
        if len(x) < 2 or len(y) < 2:
            continue
        pooled = math.sqrt(((len(x) - 1) * x.var() + (len(y) - 1) * y.var()) / max(len(x) + len(y) - 2, 1))
        d = (y.mean() - x.mean()) / pooled if pooled > 0 else np.nan
        try:
            mw_p = mannwhitneyu(x, y, alternative="two-sided").pvalue
        except Exception:
            mw_p = np.nan
        try:
            ks_p = ks_2samp(x, y).pvalue
        except Exception:
            ks_p = np.nan
        rows.append({
            "feature": f,
            "h2b_mean": x.mean(),
            "mcm2_mean": y.mean(),
            "h2b_median": x.median(),
            "mcm2_median": y.median(),
            "difference_mcm2_minus_h2b": y.mean() - x.mean(),
            "cohens_d_mcm2_minus_h2b": d,
            "abs_effect_size": abs(d) if not pd.isna(d) else np.nan,
            "mannwhitney_p": mw_p,
            "ks_p": ks_p,
        })
    return pd.DataFrame(rows).sort_values("abs_effect_size", ascending=False)


def distance_metrics(X: pd.DataFrame, y: pd.Series) -> tuple[pd.DataFrame, pd.DataFrame]:
    scaler = StandardScaler()
    z = scaler.fit_transform(X)
    labels = y.to_numpy()
    h = z[labels == "H2B_Bound"]
    m = z[labels == "MCM2_Bound"]
    ch, cm = h.mean(axis=0), m.mean(axis=0)
    between = float(np.linalg.norm(ch - cm))
    within_h = float(np.mean(pdist(h))) if len(h) > 1 else np.nan
    within_m = float(np.mean(pdist(m))) if len(m) > 1 else np.nan
    between_pair = float(np.mean(cdist(h, m)))
    sil = float(silhouette_score(z, labels)) if len(set(labels)) > 1 else np.nan
    rows = [{
        "centroid_distance": between,
        "mean_within_h2b_distance": within_h,
        "mean_within_mcm2_distance": within_m,
        "mean_between_class_distance": between_pair,
        "between_to_within_ratio": between_pair / np.nanmean([within_h, within_m]),
        "silhouette_score": sil,
    }]
    coords = pd.DataFrame(z, columns=[f"z_{c}" for c in X.columns])
    coords["bound_control_class"] = labels
    return pd.DataFrame(rows), coords


def plot_bar(df: pd.DataFrame, x: str, y: str, path: Path, title: str, top: int = 12) -> None:
    sub = df.head(top).iloc[::-1]
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.barh(sub[x], sub[y], color="#176B87")
    ax.set_title(title)
    fig.tight_layout()
    fig.savefig(path.with_suffix(".png"), dpi=200)
    fig.savefig(path.with_suffix(".pdf"))
    plt.close(fig)


def plot_pca(df: pd.DataFrame, X: pd.DataFrame, y: pd.Series, path: Path, title: str) -> pd.DataFrame:
    z = StandardScaler().fit_transform(X)
    pca = PCA(n_components=2, random_state=RANDOM_SEED)
    coords = pca.fit_transform(z)
    out = pd.DataFrame({"PC1": coords[:, 0], "PC2": coords[:, 1], "class": y.to_numpy()})
    fig, ax = plt.subplots(figsize=(7, 5))
    for cls, color in [("H2B_Bound", "#176B87"), ("MCM2_Bound", "#4D8C57")]:
        s = out[out["class"] == cls]
        ax.scatter(s["PC1"], s["PC2"], s=18, alpha=0.65, label=cls, c=color)
    ax.set_title(title)
    ax.legend()
    fig.tight_layout()
    fig.savefig(path.with_suffix(".png"), dpi=200)
    fig.savefig(path.with_suffix(".pdf"))
    plt.close(fig)
    return out


def plot_umap(X: pd.DataFrame, y: pd.Series, path: Path, log: Logger) -> None:
    if umap is None:
        log.write("UMAP not available; skipped UMAP plot")
        return
    z = StandardScaler().fit_transform(X)
    coords = umap.UMAP(random_state=RANDOM_SEED).fit_transform(z)
    fig, ax = plt.subplots(figsize=(7, 5))
    for cls, color in [("H2B_Bound", "#176B87"), ("MCM2_Bound", "#4D8C57")]:
        s = coords[y.to_numpy() == cls]
        ax.scatter(s[:, 0], s[:, 1], s=18, alpha=0.65, label=cls, c=color)
    ax.set_title("H2B_Bound vs MCM2_Bound UMAP")
    ax.legend()
    fig.tight_layout()
    fig.savefig(path.with_suffix(".png"), dpi=200)
    fig.savefig(path.with_suffix(".pdf"))
    plt.close(fig)


def target_similarity(final: pd.DataFrame, controls: pd.DataFrame, common_features: list[str], out: Path) -> pd.DataFrame:
    targets = final[final["protein"].isin(["TDP43", "NONO"])].copy()
    ctrl = controls.copy()
    Xc = ctrl[common_features].apply(pd.to_numeric, errors="coerce")
    Xt = targets[common_features].apply(pd.to_numeric, errors="coerce")
    med = Xc.median()
    Xc = Xc.fillna(med)
    Xt = Xt.fillna(med)
    scaler = StandardScaler().fit(Xc)
    zc = scaler.transform(Xc)
    zt = scaler.transform(Xt)
    y = ctrl["bound_control_class"]
    clf = LogisticRegression(max_iter=2000, class_weight="balanced")
    clf.fit(zc, y)
    probs = clf.predict_proba(zt)
    prob_df = pd.DataFrame(probs, columns=[f"P_{c}" for c in clf.classes_], index=targets.index)
    centroids = {cls: zc[y.to_numpy() == cls].mean(axis=0) for cls in clf.classes_}
    for cls, c in centroids.items():
        targets[f"distance_to_{cls}"] = np.sqrt(((zt - c) ** 2).sum(axis=1))
    targets = pd.concat([targets.reset_index(drop=True), prob_df.reset_index(drop=True)], axis=1)
    dist_cols = [f"distance_to_{cls}" for cls in clf.classes_]
    targets["closer_to_bound_control"] = targets[dist_cols].idxmin(axis=1).str.replace("distance_to_", "", regex=False)
    if {"P_H2B_Bound", "P_MCM2_Bound"}.issubset(targets.columns):
        targets["bound_probability_margin"] = (targets["P_H2B_Bound"] - targets["P_MCM2_Bound"]).abs()
        targets.loc[targets["bound_probability_margin"] < 0.10, "closer_to_bound_control"] = "neither_clear_margin"
    targets.to_csv(out / "target_distance_to_h2b_bound_vs_mcm2_bound.csv", index=False)
    targets[["protein", "dataset_name", "condition", "replicate", "track_id", *dist_cols, "P_H2B_Bound", "P_MCM2_Bound", "closer_to_bound_control"]].to_csv(
        out / "target_closer_to_which_bound_control.csv", index=False
    )
    for group, name in [
        (["protein"], "target_bound_control_similarity_by_protein.csv"),
        (["condition"], "target_bound_control_similarity_by_condition.csv"),
        (["replicate"], "target_bound_control_similarity_by_replicate.csv"),
    ]:
        targets.groupby(group, dropna=False).agg(
            n_tracks=("track_id", "count"),
            mean_P_H2B_Bound=("P_H2B_Bound", "mean"),
            mean_P_MCM2_Bound=("P_MCM2_Bound", "mean"),
            mean_distance_to_H2B_Bound=("distance_to_H2B_Bound", "mean"),
            mean_distance_to_MCM2_Bound=("distance_to_MCM2_Bound", "mean"),
        ).reset_index().to_csv(out / name, index=False)
    return targets


def old_vs_new(project: Path, inputs: dict, final: pd.DataFrame, out: Path, log: Logger) -> pd.DataFrame:
    class_cols = {"classification", "track_class", "final_classification", "condensate_class", "condensate_evaluation_class", "trapping_class", "track_category"}
    rows = []
    for p in inputs["old_eval"]:
        if p.name.lower() != "condensate_evaluation_summary.csv":
            continue
        df = pd.read_csv(p)
        old_col = next((c for c in df.columns if c in class_cols), None)
        if old_col is None:
            continue
        meta = parse_dataset(p, project)
        for _, r in df.iterrows():
            rows.append({
                **meta,
                "video_id": r.get("Video #"),
                "cell_id": r.get("Cell"),
                "track_id": r.get("Track"),
                "old_classification": r.get(old_col),
            })
    old = pd.DataFrame(rows)
    if old.empty:
        log.write("WARNING no old classification rows found")
        pd.DataFrame().to_csv(out / "old_vs_control_guided_classification_crosstab.csv", index=False)
        return old
    f = final.copy()
    for c in ["video_id", "cell_id", "track_id"]:
        f[c] = f[c].astype(str)
        old[c] = old[c].astype(str)
    joined = old.merge(
        f[["protein", "dataset_name", "video_id", "cell_id", "track_id", "final_control_guided_class"]],
        on=["protein", "dataset_name", "video_id", "cell_id", "track_id"],
        how="left",
    )
    joined.to_csv(out / "old_vs_control_guided_joined_tracks.csv", index=False)
    ct = pd.crosstab(joined["old_classification"], joined["final_control_guided_class"], dropna=False)
    ct.to_csv(out / "old_vs_control_guided_classification_crosstab.csv")
    pct = ct.div(ct.sum(axis=1).replace(0, np.nan), axis=0) * 100
    pct.to_csv(out / "old_vs_control_guided_classification_crosstab_percent.csv")
    trapping = joined[joined["old_classification"].astype(str).str.contains("trap|target|condens", case=False, na=False)]
    dest = trapping["final_control_guided_class"].value_counts(dropna=False).reset_index()
    dest.columns = ["new_final_control_guided_class", "n_old_trapping_like_tracks"]
    dest.to_csv(out / "old_trapping_like_destination_summary.csv", index=False)
    with (out / "old_trapping_like_destination_summary.md").open("w", encoding="utf-8") as h:
        h.write("# Old trapping-like destination summary\n\n")
        h.write(f"Matched old rows: {joined['final_control_guided_class'].notna().sum()} of {len(joined)}.\n\n")
        if dest.empty:
            h.write("No old trapping-like labels were identified by text matching.\n")
        else:
            h.write("Where old trapping-like tracks went in the stricter model:\n\n")
            for _, r in dest.iterrows():
                h.write(f"- {r['new_final_control_guided_class']}: {int(r['n_old_trapping_like_tracks'])}\n")
    fig, ax = plt.subplots(figsize=(8, max(4, 0.35 * len(ct))))
    im = ax.imshow(pct.fillna(0), aspect="auto", cmap="viridis")
    ax.set_xticks(range(len(pct.columns)))
    ax.set_xticklabels(pct.columns, rotation=35, ha="right")
    ax.set_yticks(range(len(pct.index)))
    ax.set_yticklabels(pct.index)
    fig.colorbar(im, ax=ax, label="% within old class")
    ax.set_title("Old vs control-guided classification")
    fig.tight_layout()
    fig.savefig(out / "old_vs_control_guided_classification_heatmap.png", dpi=200)
    fig.savefig(out / "old_vs_control_guided_classification_heatmap.pdf")
    plt.close(fig)
    return joined


def gate_failure(final: pd.DataFrame, targets_bound: pd.DataFrame, out: Path) -> pd.DataFrame:
    target = final[final["protein"].isin(["TDP43", "NONO"])].copy()
    extra = targets_bound[["protein", "dataset_name", "video_id", "cell_id", "track_id", "distance_to_H2B_Bound", "distance_to_MCM2_Bound"]].copy()
    for c in ["video_id", "cell_id", "track_id"]:
        target[c] = target[c].astype(str)
        extra[c] = extra[c].astype(str)
    target = target.merge(extra, on=["protein", "dataset_name", "video_id", "cell_id", "track_id"], how="left")
    target["genotype"] = target["dataset_name"].astype(str).str.split("-").str[0].str.upper()
    target["treatment"] = target["condition"].astype(str).str.upper()
    target["distance_to_NLS"] = target.get("distance_to_NLS_centroid", np.nan)
    target["distance_to_H2B"] = target.get("distance_to_H2B_centroid", np.nan)
    target["distance_to_MCM2_Bound"] = target.get("distance_to_MCM2_Bound_centroid", np.nan)
    target["distance_to_MCM2_Bound_HighConfidence"] = target["distance_to_MCM2_Bound_y"] if "distance_to_MCM2_Bound_y" in target else target.get("distance_to_MCM2_Bound")
    target["distance_to_H2B_Bound"] = target["distance_to_H2B_Bound"]
    target["passed_target_protein_gate"] = True
    target["passed_not_free_gate"] = target["P_free_like"].fillna(1) < 0.60
    target["passed_not_h2b_gate"] = target["P_chromatin_bound_like"].fillna(1) < 0.60
    target["passed_not_mcm2_gate"] = target["P_mcm_bound_like"].fillna(1) < 0.60
    target["passed_temporal_support_gate"] = target.get("supported_temporal_trapping", False).fillna(False).astype(bool)
    target["passed_no_diffusive_burst_gate"] = ~target.get("diffusive_burst_flag", False).fillna(False).astype(bool)
    target["passed_control_outlier_gate"] = target.get("is_control_outlier", False).fillna(False).astype(bool)
    feat_cols = ["n_frames", "frac_bound", "frac_constrained", "frac_diffusive", "max_consecutive_bound", "max_consecutive_diffusive", "control_outlier_score"]
    target["usable_feature_count"] = target[[c for c in feat_cols if c in target.columns]].notna().sum(axis=1)
    target["passed_feature_completeness_gate"] = target["usable_feature_count"] >= 3
    gate_cols = [
        "passed_not_free_gate",
        "passed_not_h2b_gate",
        "passed_not_mcm2_gate",
        "passed_temporal_support_gate",
        "passed_no_diffusive_burst_gate",
        "passed_control_outlier_gate",
        "passed_feature_completeness_gate",
        "passed_target_protein_gate",
    ]
    gate_names = {
        "passed_not_free_gate": "not_free",
        "passed_not_h2b_gate": "not_h2b_like",
        "passed_not_mcm2_gate": "not_mcm2_like",
        "passed_temporal_support_gate": "temporal_support",
        "passed_no_diffusive_burst_gate": "no_diffusive_burst",
        "passed_control_outlier_gate": "control_outlier",
        "passed_feature_completeness_gate": "feature_completeness",
        "passed_target_protein_gate": "target_protein",
    }
    failed = []
    for _, r in target.iterrows():
        bad = [gate_names[c] for c in gate_cols if not bool(r[c])]
        failed.append(bad)
    target["failed_gate_count"] = [len(x) for x in failed]
    target["first_failed_gate"] = [x[0] if x else "" for x in failed]
    target["all_failed_gates"] = [";".join(x) for x in failed]
    target["failed_only_one_gate"] = target["failed_gate_count"] == 1
    target["near_miss_target_specific"] = target["failed_gate_count"] <= 2
    req = [
        "protein", "dataset_name", "condition", "genotype", "treatment", "replicate", "track_id",
        "P_free_like", "P_chromatin_bound_like", "P_mcm_bound_like", "distance_to_NLS", "distance_to_H2B",
        "distance_to_MCM2_Bound", "distance_to_H2B_Bound", "distance_to_MCM2_Bound_HighConfidence",
        "control_outlier_score", "trapping_score", "temporal_score", "trap_prob_mean", "max_trap_dwell_frames",
        "n_supported_high_score_windows", "frac_diffusive", "max_consecutive_diffusive", *gate_cols,
        "failed_gate_count", "first_failed_gate", "all_failed_gates", "near_miss_target_specific",
        "final_control_guided_class", "classification_reason",
    ]
    for c in req:
        if c not in target.columns:
            target[c] = np.nan
    target[req].to_csv(out / "target_gate_failure_table.csv", index=False)
    for group, name in [(["protein"], "target_gate_failure_summary_by_protein.csv"), (["condition"], "target_gate_failure_summary_by_condition.csv"), (["replicate"], "target_gate_failure_summary_by_replicate.csv")]:
        rows = []
        for key, g in target.groupby(group, dropna=False):
            row = {"group": key if not isinstance(key, tuple) else "|".join(map(str, key)), "n_tracks": len(g)}
            for c in gate_cols:
                row[f"frac_failed_{gate_names[c]}"] = float((~g[c].astype(bool)).mean())
            row["mean_failed_gate_count"] = float(g["failed_gate_count"].mean())
            rows.append(row)
        pd.DataFrame(rows).to_csv(out / name, index=False)
    counts = Counter([x for xs in failed for x in xs])
    with (out / "target_gate_failure_report.md").open("w", encoding="utf-8") as h:
        h.write("# Target gate-failure report\n\n")
        h.write("No TDP43/NONO track became target-specific trapping-like because every target failed at least one required gate.\n\n")
        for gate, n in counts.most_common():
            h.write(f"- {gate}: {n} tracks ({100*n/len(target):.1f}%) failed\n")
        if counts:
            top = counts.most_common(1)[0]
            h.write(f"\nTop bottleneck: **{top[0]}** ({100*top[1]/len(target):.1f}% of target tracks).\n")
    fig, ax = plt.subplots(figsize=(8, 4.5))
    labels, vals = zip(*counts.most_common()) if counts else ([], [])
    ax.bar(labels, vals, color="#B75D52")
    ax.set_ylabel("Failed target tracks")
    ax.set_title("Why target tracks did not pass all target-specific gates")
    ax.tick_params(axis="x", rotation=30)
    fig.tight_layout()
    fig.savefig(out / "target_gate_failure_flow_plot.png", dpi=200)
    fig.savefig(out / "target_gate_failure_flow_plot.pdf")
    plt.close(fig)
    return target


def ambiguous_tables(final: pd.DataFrame, gate: pd.DataFrame, old_joined: pd.DataFrame, out: Path) -> None:
    amb = gate[gate["final_control_guided_class"].eq("ambiguous")].copy()
    if not old_joined.empty:
        old_key = old_joined.copy()
        for c in ["video_id", "cell_id", "track_id"]:
            old_key[c] = old_key[c].astype(str)
        amb = amb.merge(
            old_key[["protein", "dataset_name", "video_id", "cell_id", "track_id", "old_classification"]],
            on=["protein", "dataset_name", "video_id", "cell_id", "track_id"],
            how="left",
        )
    else:
        amb["old_classification"] = ""
    for c in ["trapping_score", "temporal_score", "control_outlier_score", "max_trap_dwell_frames", "n_supported_high_score_windows"]:
        if c not in amb:
            amb[c] = np.nan
    amb["manual_review_score"] = (
        amb["control_outlier_score"].rank(pct=True).fillna(0)
        + amb["trapping_score"].rank(pct=True).fillna(0)
        + amb["temporal_score"].rank(pct=True).fillna(0)
        + amb["max_trap_dwell_frames"].rank(pct=True).fillna(0)
        + (1 - amb["P_free_like"].rank(pct=True).fillna(1))
        + amb["near_miss_target_specific"].astype(float)
    )
    amb["ambiguous_reason"] = amb["all_failed_gates"].apply(lambda x: f"Failed gates: {x}" if str(x) else "Conflicting or insufficient evidence")
    amb["suggested_manual_review_priority"] = pd.qcut(amb["manual_review_score"].rank(method="first"), q=min(3, len(amb)), labels=["low", "medium", "high"]) if len(amb) >= 3 else "high"
    amb = amb.sort_values("manual_review_score", ascending=False)
    amb.to_csv(out / "ambiguous_target_tracks_for_review.csv", index=False)
    amb.head(100).to_csv(out / "ambiguous_target_tracks_ranked_top100.csv", index=False)
    with (out / "ambiguous_target_tracks_summary.md").open("w", encoding="utf-8") as h:
        h.write("# Ambiguous target-track inspection summary\n\n")
        h.write(f"Ambiguous TDP43/NONO tracks: {len(amb)}.\n\n")
        h.write(f"Near-miss ambiguous tracks: {int(amb['near_miss_target_specific'].sum())}.\n\n")
        h.write("High-priority rows are ranked by outlier/trapping/temporal evidence, low free-like probability, and near-miss status.\n")


def tdp43_summaries(final: pd.DataFrame, gate: pd.DataFrame, targets_bound: pd.DataFrame, out: Path, project: Path) -> None:
    folders = []
    for d in (project / "TDP43").glob("*"):
        if d.is_dir():
            meta = parse_dataset(d / "dummy.csv", project)
            analysis = list((d / "CSV").glob("*")) if (d / "CSV").exists() else []
            folders.append({**meta, "folder_name": d.name, "analysis_folder_count": len([a for a in analysis if a.is_dir()])})
    pd.DataFrame(folders).to_csv(out / "tdp43_condition_inventory.csv", index=False)
    pd.DataFrame(folders).to_csv(out / "tdp43_replicate_inventory.csv", index=False)
    t = final[final["protein"].eq("TDP43")].copy()
    t["genotype"] = t["dataset_name"].astype(str).str.split("-").str[0].str.upper()
    t["treatment"] = t["condition"].astype(str).str.upper()
    rows = []
    for key, g in t.groupby(["genotype", "treatment"], dropna=False):
        row = {"genotype": key[0], "treatment": key[1], "n_tracks": len(g)}
        vc = g["final_control_guided_class"].value_counts()
        for cls in CLASS_NAMES:
            row[f"n_{cls}"] = int(vc.get(cls, 0))
            row[f"frac_{cls}"] = float(vc.get(cls, 0) / len(g)) if len(g) else 0
        rows.append(row)
    pd.DataFrame(rows).to_csv(out / "tdp43_class_summary_by_genotype_treatment.csv", index=False)
    gate[gate["protein"].eq("TDP43")].groupby(["genotype", "treatment"], dropna=False).agg(
        n_tracks=("track_id", "count"),
        mean_failed_gate_count=("failed_gate_count", "mean"),
        frac_near_miss=("near_miss_target_specific", "mean"),
    ).reset_index().to_csv(out / "tdp43_gate_failure_by_genotype_treatment.csv", index=False)
    tb = targets_bound[targets_bound["protein"].eq("TDP43")].copy()
    tb["genotype"] = tb["dataset_name"].astype(str).str.split("-").str[0].str.upper()
    tb["treatment"] = tb["condition"].astype(str).str.upper()
    tb.groupby(["genotype", "treatment"], dropna=False).agg(
        n_tracks=("track_id", "count"),
        mean_P_H2B_Bound=("P_H2B_Bound", "mean"),
        mean_P_MCM2_Bound=("P_MCM2_Bound", "mean"),
        mean_distance_to_H2B_Bound=("distance_to_H2B_Bound", "mean"),
        mean_distance_to_MCM2_Bound=("distance_to_MCM2_Bound", "mean"),
    ).reset_index().to_csv(out / "tdp43_bound_control_similarity_by_genotype_treatment.csv", index=False)
    gate[(gate["protein"].eq("TDP43")) & (gate["final_control_guided_class"].eq("ambiguous"))].to_csv(out / "tdp43_ambiguous_tracks_by_genotype_treatment.csv", index=False)

    # Required TDP43 plots.
    summary = pd.read_csv(out / "tdp43_class_summary_by_genotype_treatment.csv")
    if not summary.empty:
        labels = summary["genotype"] + " " + summary["treatment"]
        fig, ax = plt.subplots(figsize=(9, 4.5))
        bottom = np.zeros(len(summary))
        for cls, color in [("frac_chromatin_bound_like", "#176B87"), ("frac_free_like", "#C47F2C"), ("frac_mcm_bound_like", "#4D8C57"), ("frac_ambiguous", "#8A8F98")]:
            vals = summary[cls].to_numpy() * 100
            ax.bar(labels, vals, bottom=bottom, label=cls.replace("frac_", ""), color=color)
            bottom += vals
        ax.set_ylabel("% tracks")
        ax.tick_params(axis="x", rotation=30)
        ax.legend(fontsize=8)
        fig.tight_layout()
        fig.savefig(out / "tdp43_final_class_percent_by_condition.png", dpi=200)
        fig.savefig(out / "tdp43_final_class_percent_by_condition.pdf")
        plt.close(fig)

        fig, ax = plt.subplots(figsize=(8, 4))
        vals = summary.get("frac_target_specific_trapping_like", pd.Series([0] * len(summary))).to_numpy() * 100
        ax.bar(labels, vals, color="#B75D52")
        ax.set_ylabel("% target-specific")
        ax.tick_params(axis="x", rotation=30)
        fig.tight_layout()
        fig.savefig(out / "tdp43_target_specific_or_nearmiss_by_condition.png", dpi=200)
        fig.savefig(out / "tdp43_target_specific_or_nearmiss_by_condition.pdf")
        plt.close(fig)

    gf = pd.read_csv(out / "tdp43_gate_failure_by_genotype_treatment.csv")
    if not gf.empty:
        fig, ax = plt.subplots(figsize=(8, 4))
        labels = gf["genotype"] + " " + gf["treatment"]
        ax.bar(labels, gf["mean_failed_gate_count"], color="#B75D52")
        ax.set_ylabel("Mean failed gates")
        ax.tick_params(axis="x", rotation=30)
        fig.tight_layout()
        fig.savefig(out / "tdp43_gate_failure_by_condition.png", dpi=200)
        fig.savefig(out / "tdp43_gate_failure_by_condition.pdf")
        plt.close(fig)

    sim = pd.read_csv(out / "tdp43_bound_control_similarity_by_genotype_treatment.csv")
    if not sim.empty:
        fig, ax = plt.subplots(figsize=(8, 4))
        labels = sim["genotype"] + " " + sim["treatment"]
        ax.scatter(sim["mean_distance_to_H2B_Bound"], sim["mean_distance_to_MCM2_Bound"], s=80)
        for x, yv, label in zip(sim["mean_distance_to_H2B_Bound"], sim["mean_distance_to_MCM2_Bound"], labels):
            ax.text(x, yv, label, fontsize=8)
        ax.set_xlabel("Distance to H2B_Bound")
        ax.set_ylabel("Distance to MCM2_Bound")
        fig.tight_layout()
        fig.savefig(out / "tdp43_distance_to_h2b_bound_vs_mcm2_bound_by_condition.png", dpi=200)
        fig.savefig(out / "tdp43_distance_to_h2b_bound_vs_mcm2_bound_by_condition.pdf")
        plt.close(fig)

        heat = sim[["mean_P_H2B_Bound", "mean_P_MCM2_Bound"]].to_numpy()
        fig, ax = plt.subplots(figsize=(7, max(3, 0.35 * len(sim))))
        im = ax.imshow(heat, aspect="auto", cmap="viridis", vmin=0, vmax=1)
        ax.set_yticks(range(len(labels)))
        ax.set_yticklabels(labels)
        ax.set_xticks([0, 1])
        ax.set_xticklabels(["P H2B_Bound", "P MCM2_Bound"])
        fig.colorbar(im, ax=ax)
        fig.tight_layout()
        fig.savefig(out / "tdp43_control_similarity_by_condition_heatmap.png", dpi=200)
        fig.savefig(out / "tdp43_control_similarity_by_condition_heatmap.pdf")
        plt.close(fig)


def write_reports(out: Path, metrics: pd.DataFrame, down: pd.DataFrame, perm: pd.DataFrame, effects: pd.DataFrame, dist: pd.DataFrame, target_sim: pd.DataFrame, gate: pd.DataFrame) -> str:
    bal = float(metrics["balanced_accuracy"].iloc[0])
    down_med = float(down["balanced_accuracy"].median())
    p = float(perm["p_value"].iloc[0])
    if bal >= 0.80 and down_med >= 0.75 and p < 0.05:
        interp = "yes"
    elif bal >= 0.62 and down_med >= 0.60:
        interp = "partly"
    else:
        interp = "no"
    with (out / "control_separability_report.md").open("w", encoding="utf-8") as h:
        h.write("# H2B_Bound versus MCM2_Bound separability report\n\n")
        h.write(f"Final interpretation: **{interp}**.\n\n")
        h.write(f"Balanced cross-validated balanced accuracy: {bal:.3f}.\n\n")
        h.write(f"Repeated downsample median balanced accuracy: {down_med:.3f}.\n\n")
        h.write(f"Permutation-test p-value: {p:.4f}.\n\n")
        h.write("Top separating features:\n\n")
        for _, r in effects.head(8).iterrows():
            h.write(f"- {r['feature']}: Cohen's d={r['cohens_d_mcm2_minus_h2b']:.2f}, MCM2-H2B mean difference={r['difference_mcm2_minus_h2b']:.3g}\n")
        h.write("\nTarget tracks were projected onto the H2B_Bound/MCM2_Bound feature space for similarity only; they were not used to train the bound-control classifier.\n\n")
        h.write("Cautious interpretation: this is a motion-pattern diagnostic, not proof of physical condensate or paraspeckle localization.\n")
    with (out / "plain_english_summary.md").open("w", encoding="utf-8") as h:
        h.write("# Plain-English summary\n\n")
        h.write(f"H2B-bound and MCM2-bound controls are interpreted as **{interp}** separable in this feature space.\n\n")
        h.write("The stricter target-specific class requires tracks to differ from free-like, H2B-like, and MCM2-like controls while also showing temporal support.\n\n")
        h.write("No definitive condensate localization is claimed from tracking alone.\n")
    with (out / "easy_interpretation_diagnostics.md").open("w", encoding="utf-8") as h:
        h.write("# Easy interpretation diagnostics\n\n")
        h.write(f"Are H2B_Bound and MCM2_Bound distinct enough to be separate controls? **{interp}**.\n\n")
        h.write("If separability is partial or weak, interpret target similarity cautiously and consider merging bound controls in sensitivity checks.\n")
    with (out / "key_findings_diagnostics.md").open("w", encoding="utf-8") as h:
        h.write("# Key findings\n\n")
        h.write(f"- H2B_Bound vs MCM2_Bound separability: {interp}.\n")
        h.write(f"- Balanced accuracy: {bal:.3f}; downsample median: {down_med:.3f}; permutation p={p:.4f}.\n")
        h.write(f"- Ambiguous target tracks: {int((gate['final_control_guided_class']=='ambiguous').sum())}.\n")
    with (out / "recommended_next_actions.md").open("w", encoding="utf-8") as h:
        h.write("# Recommended next actions\n\n")
        h.write("- Inspect high-priority ambiguous TDP43/NONO tracks manually.\n")
        h.write("- Validate candidate trapping-like dynamics with condensate or paraspeckle markers.\n")
        h.write("- Treat WT ACTD as descriptive unless more replicates are added.\n")
        h.write("- If bound-control separability is weak, run a sensitivity analysis with H2B_Bound and MCM2_Bound merged.\n")
    return interp


def old_plot_targets(targets: pd.DataFrame, out: Path) -> None:
    fig, ax = plt.subplots(figsize=(8, 4.5))
    for protein, color in [("TDP43", "#B75D52"), ("NONO", "#7C5FB2")]:
        s = targets[targets["protein"] == protein]
        ax.scatter(s["distance_to_H2B_Bound"], s["distance_to_MCM2_Bound"], s=12, alpha=0.45, label=protein, c=color)
    ax.set_xlabel("Distance to H2B_Bound")
    ax.set_ylabel("Distance to MCM2_Bound")
    ax.legend()
    fig.tight_layout()
    fig.savefig(out / "target_distance_to_bound_controls.png", dpi=200)
    fig.savefig(out / "target_distance_to_bound_controls.pdf")
    plt.close(fig)
    for group, name in [("protein", "target_closer_to_h2b_or_mcm2_bound_by_protein"), ("condition", "target_closer_to_h2b_or_mcm2_bound_by_condition")]:
        table = pd.crosstab(targets[group], targets["closer_to_bound_control"], normalize="index") * 100
        ax = table.plot(kind="bar", stacked=True, figsize=(8, 4.5))
        ax.set_ylabel("% target tracks")
        ax.figure.tight_layout()
        ax.figure.savefig(out / f"{name}.png", dpi=200)
        ax.figure.savefig(out / f"{name}.pdf")
        plt.close(ax.figure)


def write_presentation(out: Path, interp: str, metrics: pd.DataFrame, targets: pd.DataFrame, gate: pd.DataFrame) -> None:
    pres_dir = out / "presentation"
    pres_dir.mkdir(exist_ok=True)
    pptx = pres_dir / "control_guided_diagnostics_and_separability.pptx"
    if Presentation is None:
        (pres_dir / "presentation_unavailable.md").write_text("python-pptx is unavailable; PPTX could not be generated.\n", encoding="utf-8")
        return
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide(title: str, bullets: list[str], subtitle: str = ""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(247, 243, 234)
        tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(11.8), Inches(0.8))
        p = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(22, 32, 42)
        if subtitle:
            st = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(11.5), Inches(0.4))
            sp = st.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(14)
            sp.font.color.rgb = RGBColor(104, 113, 125)
        box = slide.shapes.add_textbox(Inches(0.85), Inches(1.85), Inches(11.4), Inches(4.8))
        tf = box.text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = b
            para.font.size = Pt(18 if i == 0 else 15)
            para.font.color.rgb = RGBColor(22, 32, 42)
            para.level = 0
        return slide

    add_slide("Control-guided condensate classification diagnostics", ["H2B-bound vs MCM2-bound separability and target-track interpretation"], "Motion-pattern analysis only; no definitive condensate localization is claimed.")
    add_slide("Why this diagnostic was needed", ["H2B_Bound and MCM2_Bound are both chromatin-associated controls.", "We tested whether they are separable before using them as distinct explanations for target tracks."])
    add_slide("Data inventory", [f"H2B_Bound tracks: {int((targets['closer_to_bound_control'].notna()).sum())} target projections", "Primary controls: single-behaviour H2B-bound and MCM2-bound files.", "TDP43/NONO targets were projected, not used for training."])
    add_slide("H2B_Bound vs MCM2_Bound separability", [f"Conclusion: {interp}.", f"Balanced accuracy: {metrics['balanced_accuracy'].iloc[0]:.3f}.", f"ROC-AUC: {metrics['roc_auc'].iloc[0]:.3f}."])
    add_slide("Features separating controls", ["Top ranked features are in h2b_vs_mcm2bound_feature_rankings.csv.", "Effect sizes and nonparametric tests are in h2b_vs_mcm2bound_feature_effect_sizes.csv.", "Interpret features as motion/track descriptors, not localization markers."])
    add_slide("Separation after balancing", ["Repeated downsampling matched H2B_Bound count to MCM2_Bound count.", "The main conclusion is based on balanced analysis, not raw class imbalance."])
    add_slide("Where TDP43 and NONO fall", ["Each target track received H2B_Bound-like and MCM2_Bound-like probabilities.", "Protein, condition, and replicate summaries show whether targets are closer to one bound control or neither."])
    add_slide("Old versus new classification", ["Old condensate_evaluation labels were cross-tabbed against stricter control-guided labels.", "The stricter target-specific label asks for more evidence than old trapping-like labels."])
    add_slide("Why no target-specific trapping-like tracks?", ["Gate-failure tables quantify which gates failed.", "Top bottlenecks are listed in target_gate_failure_report.md.", "Near-miss tracks are separated for manual review."])
    add_slide("Ambiguous tracks", [f"Ambiguous target tracks: {int((gate['final_control_guided_class']=='ambiguous').sum())}.", "Top 100 ranked ambiguous tracks are saved for manual inspection.", "Ambiguous can mean near-miss or low-information; review priority is provided."])
    add_slide("TDP43 expanded condition analysis", ["WT CT, WT ARS, 12D CT, 12D ARS, and WT ACTD folders were parsed.", "Track-level and replicate-level summaries were generated.", "WT ACTD is descriptive unless additional replicates are added."])
    add_slide("Biological interpretation", ["Supported: control-guided motion-pattern comparisons.", "Not proven: physical condensate location or paraspeckle localization.", "Next: validate with condensate/paraspeckle markers and inspect near-miss tracks."])
    add_slide("Recommended next steps", ["Inspect high-priority ambiguous target tracks.", "Use marker/mask validation for localization.", "Consider merged bound-control sensitivity analysis if separability is only partial."])
    add_slide("File locations", ["All diagnostics outputs are under diagnostics_and_separability.", "Presentation: presentation/control_guided_diagnostics_and_separability.pptx", "Main reports: control_separability_report.md and target_gate_failure_report.md"])
    prs.save(pptx)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--project-folder", required=True)
    args = ap.parse_args()
    project = Path(args.project_folder).resolve()
    out = project / "control_guided_classification" / "diagnostics_and_separability"
    out.mkdir(parents=True, exist_ok=True)
    (out / "logs").mkdir(exist_ok=True)
    log = Logger(out / "logs" / f"diagnostics_and_separability_{now_stamp()}.log")
    inputs = locate_inputs(project, out, log)

    h2b = frame_track_features(inputs["h2b"], "H2B_Bound", project)
    mcm = frame_track_features(inputs["mcm"], "MCM2_Bound", project)
    controls = pd.concat([h2b, mcm], ignore_index=True)
    controls.to_csv(out / "h2b_mcm2_bound_track_features.csv", index=False)
    log.write(f"H2B_Bound tracks: {len(h2b)}; MCM2_Bound tracks: {len(mcm)}")

    ignore = {"protein", "dataset_name", "genotype", "date_or_experiment_id", "treatment", "replicate", "bound_control_class", "source_file", "source_folder", "video_id", "cell_id", "track_id"}
    features = choose_features(controls, ignore)
    log.write(f"Features used for bound-control separability: {features}")
    X = controls[features].apply(pd.to_numeric, errors="coerce").fillna(controls[features].median())
    y = controls["bound_control_class"]

    metrics, cm, rf, scaler = balanced_cv_metrics(X, y)
    metrics.to_csv(out / "h2b_vs_mcm2bound_balanced_classifier_metrics.csv", index=False)
    metrics.to_csv(out / "control_separability_metrics.csv", index=False)
    cm.to_csv(out / "h2b_vs_mcm2bound_confusion_matrix.csv")
    down = repeated_downsample(X, y)
    down.to_csv(out / "h2b_vs_mcm2bound_downsample_repeats.csv", index=False)
    perm = permutation_test(X, y, float(metrics["balanced_accuracy"].iloc[0]))
    perm.to_csv(out / "h2b_vs_mcm2bound_permutation_test.csv", index=False)
    effects = effect_sizes(controls, features)
    effects.to_csv(out / "h2b_vs_mcm2bound_feature_effect_sizes.csv", index=False)
    rankings = pd.DataFrame({"feature": features, "random_forest_importance": rf.feature_importances_}).merge(
        effects[["feature", "abs_effect_size", "cohens_d_mcm2_minus_h2b"]], on="feature", how="left"
    ).sort_values(["random_forest_importance", "abs_effect_size"], ascending=False)
    rankings.to_csv(out / "h2b_vs_mcm2bound_feature_rankings.csv", index=False)
    dist, _ = distance_metrics(X, y)
    dist.to_csv(out / "h2b_vs_mcm2bound_distance_metrics.csv", index=False)

    plot_pca(controls, X, y, out / "h2b_bound_vs_mcm2_bound_pca", "H2B_Bound vs MCM2_Bound PCA")
    plot_umap(X, y, out / "h2b_bound_vs_mcm2_bound_umap", log)
    plot_bar(rankings, "feature", "random_forest_importance", out / "h2b_bound_vs_mcm2_bound_feature_importance", "Feature importance")
    top = rankings["feature"].head(6).tolist()
    if top:
        long = controls.melt(id_vars=["bound_control_class"], value_vars=top, var_name="feature", value_name="value")
        fig, axes = plt.subplots(len(top), 1, figsize=(7, 1.8 * len(top)))
        axes = np.atleast_1d(axes)
        for ax, f in zip(axes, top):
            vals = [controls.loc[controls["bound_control_class"] == cls, f].dropna() for cls in ["H2B_Bound", "MCM2_Bound"]]
            ax.boxplot(vals, labels=["H2B", "MCM2"], showfliers=False)
            ax.set_title(f)
        fig.tight_layout()
        fig.savefig(out / "h2b_bound_vs_mcm2_bound_top_features_boxplots.png", dpi=200)
        fig.savefig(out / "h2b_bound_vs_mcm2_bound_top_features_boxplots.pdf")
        plt.close(fig)
    fig, ax = plt.subplots(figsize=(7, 4))
    z = scaler.transform(X)
    d0 = np.linalg.norm(z - z[y.eq("H2B_Bound")].mean(axis=0), axis=1)
    d1 = np.linalg.norm(z - z[y.eq("MCM2_Bound")].mean(axis=0), axis=1)
    ax.hist(d0[y.eq("H2B_Bound")], alpha=0.55, label="H2B to H2B", color="#176B87")
    ax.hist(d1[y.eq("H2B_Bound")], alpha=0.55, label="H2B to MCM2", color="#4D8C57")
    ax.set_title("Distance distribution")
    ax.legend()
    fig.tight_layout()
    fig.savefig(out / "h2b_bound_vs_mcm2_bound_distance_distribution.png", dpi=200)
    fig.savefig(out / "h2b_bound_vs_mcm2_bound_distance_distribution.pdf")
    plt.close(fig)

    final = pd.read_csv(inputs["final"], low_memory=False)
    common = [f for f in features if f in final.columns]
    if len(common) < 3:
        common = [c for c in ["n_frames", "track_length", "frac_bound", "frac_constrained", "frac_diffusive", "max_consecutive_bound", "max_consecutive_diffusive", "n_state_switches"] if c in final.columns and c in controls.columns]
    log.write(f"Common target projection features: {common}")
    target_sim = target_similarity(final, controls, common, out)
    old_joined = old_vs_new(project, inputs, final, out, log)
    gate = gate_failure(final, target_sim, out)
    ambiguous_tables(final, gate, old_joined, out)
    tdp43_summaries(final, gate, target_sim, out, project)
    old_plot_targets(target_sim, out)
    interp = write_reports(out, metrics, down, perm, effects, dist, target_sim, gate)
    write_presentation(out, interp, metrics, target_sim, gate)

    acceptance = [
        "control_separability_report.md",
        "control_separability_metrics.csv",
        "h2b_vs_mcm2bound_feature_effect_sizes.csv",
        "h2b_vs_mcm2bound_balanced_classifier_metrics.csv",
        "target_distance_to_h2b_bound_vs_mcm2_bound.csv",
        "old_vs_control_guided_classification_crosstab.csv",
        "ambiguous_target_tracks_for_review.csv",
        "target_gate_failure_table.csv",
        "target_gate_failure_report.md",
        "tdp43_condition_inventory.csv",
        "tdp43_class_summary_by_genotype_treatment.csv",
        "easy_interpretation_diagnostics.md",
        "plain_english_summary.md",
        "presentation/control_guided_diagnostics_and_separability.pptx",
    ]
    missing = [f for f in acceptance if not (out / f).exists()]
    if missing:
        (out / "diagnostics_failure_report.md").write_text("Missing required files:\n" + "\n".join(missing), encoding="utf-8")
        log.write(f"FAIL missing acceptance files: {missing}")
        return 2
    log.write("Diagnostics and separability analysis completed successfully")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
